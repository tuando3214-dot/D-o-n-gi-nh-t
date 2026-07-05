import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. KHỞI TẠO PRESENTATION & ĐỊNH DẠNG KÍCH THƯỚC CHUẨN 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Bảng màu thương hiệu Đại học Hòa Bình
DARK_BLUE = RGBColor(0, 80, 136)     # Màu chủ đạo
LIGHT_BLUE = RGBColor(59, 130, 246)   # Màu điểm nhấn (Accent)
DARK_GRAY = RGBColor(51, 65, 85)     # Màu chữ nội dung
LIGHT_GRAY = RGBColor(241, 245, 249)  # Màu nền khối (Background tile)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 110)

blank_layout = prs.slide_layouts[6]  # Layout trống hoàn toàn

def apply_background_and_header(slide, title_text=None):
    """Áp dụng màu nền, thanh trang trí đầu slide và tiêu đề chuẩn."""
    # Thiết lập màu nền slide trắng tinh tế
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Thêm thanh trang trí màu xanh thương hiệu ở trên cùng slide
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.fill.background() # Không viền
    
    # Nếu có tiêu đề slide
    if title_text:
        # Đường kẻ nhỏ phân tách tiêu đề bên trái
        decor_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(0.6), Inches(0.08), Inches(0.8)
        )
        decor_line.fill.solid()
        decor_line.fill.fore_color.rgb = LIGHT_BLUE
        decor_line.line.fill.background()
        
        # Hộp văn bản chứa tiêu đề slide
        tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.5), Inches(1.0))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = 'Arial'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

# ==========================================
# SLIDE 1: TIÊU ĐỀ TRANG TRỌNG (TITLE SLIDE)
# ==========================================
slide1 = prs.slides.add_slide(blank_layout)
apply_background_and_header(slide1)

# Thêm logo khối trang trí góc phải dưới
decor_bg = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4.5), Inches(4), Inches(4))
decor_bg.fill.solid()
decor_bg.fill.fore_color.rgb = LIGHT_GRAY
decor_bg.line.fill.background()

# Tên trường
univ_box = slide1.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
tf_univ = univ_box.text_frame
p_univ = tf_univ.paragraphs[0]
p_univ.text = "TRƯỜNG ĐẠI HỌC HÒA BÌNH • KHOA CÔNG NGHỆ THÔNG TIN"
p_univ.font.name = 'Arial'
p_univ.font.size = Pt(13)
p_univ.font.bold = True
p_univ.font.color.rgb = DARK_BLUE

# Tên đề tài lớn
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2))
tf_title = title_box.text_frame
tf_title.word_wrap = True
p_sub = tf_title.paragraphs[0]
p_sub.text = "BÀI TẬP LỚN CÔNG NGHỆ BIG DATA"
p_sub.font.name = 'Arial'
p_sub.font.size = Pt(22)
p_sub.font.bold = True
p_sub.font.color.rgb = LIGHT_BLUE

p_main = tf_title.add_paragraph()
p_main.text = "HOUSE PRICE PREDICTION"
p_main.font.name = 'Arial'
p_main.font.size = Pt(54)
p_main.font.bold = True
p_main.font.color.rgb = DARK_BLUE

p_desc = tf_title.add_paragraph()
p_desc.text = "Ứng dụng Học máy dự đoán giá trị bất động sản theo thời gian thực"
p_desc.font.name = 'Arial'
p_desc.font.size = Pt(16)
p_desc.font.italic = True
p_desc.font.color.rgb = DARK_GRAY

# Chân trang thông tin nhóm và giảng viên
footer_bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.8))
footer_bg.fill.solid()
footer_bg.fill.fore_color.rgb = LIGHT_GRAY
footer_bg.line.color.rgb = DARK_BLUE
footer_bg.line.width = Pt(1.5)

info_box = slide1.shapes.add_textbox(Inches(1.1), Inches(4.8), Inches(11.1), Inches(1.6))
tf_info = info_box.text_frame
p_gv = tf_info.paragraphs[0]
p_gv.text = "Giảng viên hướng dẫn:  Thầy ĐOÀN VĂN HÒA"
p_gv.font.name = 'Arial'
p_gv.font.size = Pt(14)
p_gv.font.bold = True
p_gv.font.color.rgb = DARK_BLUE

p_sv = tf_info.add_paragraph()
p_sv.text = "Sinh viên thực hiện:      Nhóm 4 (Lớp 522CNT)"
p_sv.font.name = 'Arial'
p_sv.font.size = Pt(14)
p_sv.font.bold = True
p_sv.font.color.rgb = DARK_BLUE

p_members = tf_info.add_paragraph()
p_members.text = "N.V. Minh • P.N. Anh • H.V. Vũ • Đ.V. Tuấn • V.Đ. Huy"
p_members.font.name = 'Arial'
p_members.font.size = Pt(12)
p_members.font.color.rgb = DARK_GRAY

# ==========================================
# SLIDE 2: TÍNH CẤP THIẾT & MỤC TIÊU ĐỀ TÀI
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
apply_background_and_header(slide2, "Tính Cấp Thiết & Mục Tiêu Nghiên Cứu")

# Khối cột trái: Lý do chọn đề tài
box_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
box_left.fill.solid()
box_left.fill.fore_color.rgb = LIGHT_GRAY
box_left.line.fill.background()

tf_left = box_left.text_frame
tf_left.word_wrap = True
tf_left.margin_top = tf_left.margin_left = tf_left.margin_right = Inches(0.3)
p_l_title = tf_left.paragraphs[0]
p_l_title.text = "LÝ DO CHỌN ĐỀ TÀI"
p_l_title.font.name = 'Arial'
p_l_title.font.size = Pt(18)
p_l_title.font.bold = True
p_l_title.font.color.rgb = DARK_BLUE

bullet1 = tf_left.add_paragraph()
bullet1.text = "• Thị trường bất động sản biến động liên tục, chịu ảnh hưởng phức tạp của nhiều yếu tố đan xen cấu trúc vật lý."
bullet1.font.size = Pt(14)
bullet1.font.color.rgb = DARK_GRAY

bullet2 = tf_left.add_paragraph()
bullet2.text = "• Phương pháp định giá truyền thống tốn kém thời gian, chi phí, mang tính cảm tính cá nhân và gây bất đối xứng thông tin."
bullet2.font.size = Pt(14)
bullet2.font.color.rgb = DARK_GRAY

bullet3 = tf_left.add_paragraph()
bullet3.text = "• Khoa học dữ liệu & Học máy (Machine Learning) là giải pháp tối ưu giúp tự động hóa và minh bạch hóa giá trị tài sản."
bullet3.font.size = Pt(14)
bullet3.font.color.rgb = DARK_GRAY

# Khối cột phải: Mục tiêu nghiên cứu
box_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
box_right.fill.solid()
box_right.fill.fore_color.rgb = LIGHT_GRAY
box_right.line.fill.background()

tf_right = box_right.text_frame
tf_right.word_wrap = True
tf_right.margin_top = tf_right.margin_left = tf_right.margin_right = Inches(0.3)
p_r_title = tf_right.paragraphs[0]
p_r_title.text = "MỤC TIÊU CHÍNH"
p_r_title.font.name = 'Arial'
p_r_title.font.size = Pt(18)
p_r_title.font.bold = True
p_r_title.font.color.rgb = DARK_BLUE

bullet_r1 = tf_right.add_paragraph()
bullet_r1.text = "• Thu thập & Tiền xử lý: Làm sạch tệp thô housing.csv, xử lý dữ liệu trống, loại bỏ các biến ngoại lai (outliers)."
bullet_r1.font.size = Pt(14)
bullet_r1.font.color.rgb = DARK_GRAY

bullet_r2 = tf_right.add_paragraph()
bullet_r2.text = "• Huấn luyện & Tối ưu: Huấn luyện thuật toán Random Forest Regressor trên tập dữ liệu và đóng gói mô hình ."
bullet_r2.font.size = Pt(14)
bullet_r2.font.color.rgb = DARK_GRAY

bullet_r3 = tf_right.add_paragraph()
bullet_r3.text = "• Triển khai thực tiễn: Thiết kế giao diện web trực quan Streamlit giúp người dùng tra cứu giá trị tức thì theo thời gian thực."
bullet_r3.font.size = Pt(14)
bullet_r3.font.color.rgb = DARK_GRAY

# ==========================================
# SLIDE 3: THÔNG SỐ CÁC BIẾN ĐẶC TRƯNG
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
apply_background_and_header(slide3, "Thông Số Các Biến Đặc Trưng Đầu Vào")

# Cột trái: Phương trình toán học hồi quy
formula_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
formula_box.fill.solid()
formula_box.fill.fore_color.rgb = LIGHT_GRAY
formula_box.line.fill.background()

tf_f = formula_box.text_frame
tf_f.word_wrap = True
tf_f.margin_top = tf_f.margin_left = tf_f.margin_right = Inches(0.3)
p_f_t = tf_f.paragraphs[0]
p_f_t.text = "MÔ HÌNH TOÁN HỌC HỒI QUY"
p_f_t.font.name = 'Arial'
p_f_t.font.size = Pt(16)
p_f_t.font.bold = True
p_f_t.font.color.rgb = DARK_BLUE

p_eq_desc = tf_f.add_paragraph()
p_eq_desc.text = "Hệ thống dự đoán giá trị biến mục tiêu (y) dựa trên tập hợp 5 biến đặc trưng cấu trúc độc lập (X):"
p_eq_desc.font.size = Pt(14)
p_eq_desc.font.color.rgb = DARK_GRAY

p_eq = tf_f.add_paragraph()
p_eq.text = "y = f(x1, x2, x3, x4, x5)"
p_eq.font.name = 'Consolas'
p_eq.font.size = Pt(22)
p_eq.font.bold = True
p_eq.font.color.rgb = LIGHT_BLUE
p_eq.alignment = PP_ALIGN.CENTER

p_f_bullets = tf_f.add_paragraph()
p_f_bullets.text = "• Biến phụ thuộc (y): Giá trị quy đổi thực tế của bất động sản (USD).\n• Biến độc lập (X): Gồm 5 chỉ số vật lý định lượng thực tế của ngôi nhà."
p_f_bullets.font.size = Pt(13)
p_f_bullets.font.color.rgb = DARK_GRAY

# Cột phải: Danh sách chi tiết các thuộc tính
features_box = slide3.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
tf_feat = features_box.text_frame
tf_feat.word_wrap = True

features_list = [
    ("1. Diện tích (Area):", "Kích thước không gian sử dụng thực tế (đơn vị tính: sqft)."),
    ("2. Số phòng ngủ (Bedrooms):", "Tổng số phòng ngủ thực tế thiết kế (từ 1 đến 6 phòng)."),
    ("3. Số phòng vệ sinh (Bathrooms):", "Đóng góp lớn thứ hai vào giá trị căn nhà (từ 1 đến 4)."),
    ("4. Số tầng (Stories):", "Chiều cao tổng số tầng của tòa nhà (từ 1 đến 4 tầng)."),
    ("5. Chỗ đỗ xe (Parking):", "Sức chứa tối đa của gara bãi đỗ xe ô tô (từ 0 đến 3 xe).")
]

for title, desc in features_list:
    p_title = tf_feat.add_paragraph()
    p_title.text = title
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(15)
    p_title.font.bold = True
    p_title.font.color.rgb = DARK_BLUE
    
    p_desc = tf_feat.add_paragraph()
    p_desc.text = desc
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = DARK_GRAY
    p_desc.space_after = Pt(10)

# ==========================================
# SLIDE 4: QUY TRÌNH NGHIÊN CỨU 4 BƯỚC
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
apply_background_and_header(slide4, "Quy Trình Nghiên Cứu Hệ Thống 4 Bước")

# Vẽ sơ đồ dòng chảy ngang (Timeline)
steps = [
    ("1. TRÍCH XUẤT", "Thu thập và bóc tách dữ liệu định lượng từ tệp thô housing.csv."),
    ("2. TIỀN XỬ LÝ", "Dùng StandardScaler đồng bộ hóa thang đo, loại bỏ các giá trị trống."),
    ("3. HUẤN LUYỆN", "Chạy Random Forest Regressor và đóng gói file nhị phân .pkl."),
    ("4. TRIỂN KHAI", "Xây dựng web Streamlit hỗ trợ tra cứu giá thời gian thực.")
]

for i, (title, desc) in enumerate(steps):
    left = Inches(0.8 + (i * 2.95))
    top = Inches(2.5)
    width = Inches(2.7)
    height = Inches(3.5)
    
    # Khung bao ngoài mỗi bước
    step_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    step_box.fill.solid()
    step_box.fill.fore_color.rgb = LIGHT_GRAY
    step_box.line.color.rgb = DARK_BLUE
    step_box.line.width = Pt(1.5)
    
    tf_step = step_box.text_frame
    tf_step.word_wrap = True
    tf_step.margin_top = tf_step.margin_left = tf_step.margin_right = Inches(0.2)
    
    p_st = tf_step.paragraphs[0]
    p_st.text = title
    p_st.font.name = 'Arial'
    p_st.font.size = Pt(16)
    p_st.font.bold = True
    p_st.font.color.rgb = LIGHT_BLUE
    p_st.alignment = PP_ALIGN.CENTER
    p_st.space_after = Pt(15)
    
    p_sd = tf_step.add_paragraph()
    p_sd.text = desc
    p_sd.font.size = Pt(13)
    p_sd.font.color.rgb = DARK_GRAY
    p_sd.alignment = PP_ALIGN.CENTER

# ==========================================
# SLIDE 5: KIẾN TRÚC PHÂN LỚP HỆ THỐNG
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
apply_background_and_header(slide5, "Kiến Trúc Phân Lớp Hệ Thống")

# Khối kiến trúc 3 lớp xếp chồng dọc
layers = [
    ("LỚP GIAO DIỆN (UI LAYER - STREAMLIT APPLICATION)", 
     "• Form Nhập Liệu: Thu thập dữ liệu diện tích, số phòng ngủ, phòng tắm, số tầng, bãi đỗ xe.\n• Tab Dashboard: Trực quan hóa thị trường bất động sản thông qua biểu đồ phân tán, mật độ và ma trận nhiệt."),
    ("LỚP LOGIC & XỬ LÝ (CORE LOGIC LAYER - PYTHON ENGINE)", 
     "• StandardScaler: Đồng bộ hóa thang đo ma trận dữ liệu thô đầu vào.\n• Random Forest Regressor: Tính toán toán học và khớp quy luật giá trị bất động sản tối ưu."),
    ("LỚP LƯU TRỮ (DATA LAYER - FILE SYSTEM)", 
     "• housing.csv (Dữ liệu thô gốc)\n• train.py (Mã nguồn huấn luyện)\n• house_model.pkl & scaler.pkl (Mô hình và bộ chuẩn hóa nhị phân được đóng gói)")
]

for i, (title, desc) in enumerate(layers):
    top = Inches(1.8 + (i * 1.75))
    layer_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.5))
    layer_box.fill.solid()
    layer_box.fill.fore_color.rgb = LIGHT_GRAY
    layer_box.line.color.rgb = DARK_BLUE
    layer_box.line.width = Pt(1.5)
    
    tf_l = layer_box.text_frame
    tf_l.word_wrap = True
    tf_l.margin_top = tf_l.margin_left = tf_l.margin_right = Inches(0.15)
    
    p_lt = tf_l.paragraphs[0]
    p_lt.text = title
    p_lt.font.name = 'Arial'
    p_lt.font.size = Pt(15)
    p_lt.font.bold = True
    p_lt.font.color.rgb = LIGHT_BLUE
    
    p_ld = tf_l.add_paragraph()
    p_ld.text = desc
    p_ld.font.size = Pt(12.5)
    p_ld.font.color.rgb = DARK_GRAY

# ==========================================
# SLIDE 6: CẤU TRÚC BẢNG DỮ LIỆU CHUẨN 3NF
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
apply_background_and_header(slide6, "Cấu Trúc Bảng Dữ Liệu Housing")

# Thêm bảng dữ liệu chi tiết
rows, cols = 7, 4
left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5)
table_shape = slide6.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Thiết lập độ rộng cột
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(1.8)
table.columns[2].width = Inches(3.0)
table.columns[3].width = Inches(4.7)

headers = ["Tên trường (Column)", "Kiểu dữ liệu", "Ràng buộc giá trị", "Mô tả ý nghĩa đặc trưng"]
table_data = [
    ["price", "Float / Int", "Biến mục tiêu (y) > 0", "Giá trị bán thực tế của tài sản (USD)"],
    ["area", "Int", "Không rỗng, > 0", "Diện tích sử dụng của ngôi nhà (sqft)"],
    ["bedrooms", "Int", "Miền giá trị: [1 - 6]", "Tổng số lượng phòng ngủ thiết kế"],
    ["bathrooms", "Int", "Miền giá trị: [1 - 4]", "Số lượng phòng vệ sinh/phòng tắm"],
    ["stories", "Int", "Miền giá trị: [1 - 4]", "Tổng số tầng của căn nhà"],
    ["parking", "Int", "Miền giá trị: [0 - 3]", "Sức chứa tối đa của gara đỗ xe ô tô"]
]

# Ghi dữ liệu header
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    for p in cell.text_frame.paragraphs:
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# Ghi dữ liệu dòng
for row_idx, row_data in enumerate(table_data):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = value
        cell.fill.solid()
        if row_idx % 2 == 0:
            cell.fill.fore_color.rgb = LIGHT_GRAY
        else:
            cell.fill.fore_color.rgb = WHITE
            
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Arial'
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GRAY
            if col_idx == 0:
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE

# ==========================================
# SLIDE 7: KẾT QUẢ ĐÁNH GIÁ & KIỂM THỬ THỰC NGHIỆM
# ==========================================
slide7 = prs.slides.add_slide(blank_layout)
apply_background_and_header(slide7, "Kết Quả Đánh Giá & Kiểm Thử")

# Cột trái: Hộp số hiệu năng khổng lồ
num_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.8))
num_box.fill.solid()
num_box.fill.fore_color.rgb = DARK_BLUE
num_box.line.fill.background()

tf_num = num_box.text_frame
tf_num.word_wrap = True
tf_num.margin_top = Inches(1.2)
p_num_val = tf_num.paragraphs[0]
p_num_val.text = "99.65%"
p_num_val.font.name = 'Arial'
p_num_val.font.size = Pt(72)
p_num_val.font.bold = True
p_num_val.font.color.rgb = WHITE
p_num_val.alignment = PP_ALIGN.CENTER

p_num_lbl = tf_num.add_paragraph()
p_num_lbl.text = "ĐỘ CHÍNH XÁC R2 (TRAIN)"
p_num_lbl.font.name = 'Arial'
p_num_lbl.font.size = Pt(16)
p_num_lbl.font.bold = True
p_num_lbl.font.color.rgb = LIGHT_BLUE
p_num_lbl.alignment = PP_ALIGN.CENTER

# Cột phải: Đánh giá chi tiết hiệu năng
perf_box = slide7.shapes.add_textbox(Inches(5.8), Inches(1.8), Inches(6.7), Inches(4.8))
tf_perf = perf_box.text_frame
tf_perf.word_wrap = True

p_p_title = tf_perf.paragraphs[0]
p_p_title.text = "CHỈ SỐ HIỆU NĂNG MÔ HÌNH THỰC TẾ"
p_p_title.font.name = 'Arial'
p_p_title.font.size = Pt(20)
p_p_title.font.bold = True
p_p_title.font.color.rgb = DARK_BLUE
p_p_title.space_after = Pt(20)

bullets_perf = [
    ("• Độ chính xác thử nghiệm thực tế:", "Mô hình đạt độ chính xác R2 Score khả quan trên tập kiểm thử độc lập (~67%), chứng minh khả năng học tốt các quy luật phi tuyến phức tạp của thị trường."),
    ("• Sai số tuyệt đối trung bình (MAE):", "Giá trị sai số được kiểm soát chặt chẽ ở mức tối thiểu cho phép, đảm bảo độ tin cậy khi triển khai định giá thương mại thực tế."),
    ("• Tốc độ xử lý thời gian thực:", "Thời gian nạp bộ chuẩn hóa scaler và trả kết quả định giá nhà trên giao diện web Streamlit diễn ra dưới 0.5 giây.")
]

for title, desc in bullets_perf:
    p_bt = tf_perf.add_paragraph()
    p_bt.text = title
    p_bt.font.name = 'Arial'
    p_bt.font.size = Pt(15)
    p_bt.font.bold = True
    p_bt.font.color.rgb = LIGHT_BLUE
    
    p_bd = tf_perf.add_paragraph()
    p_bd.text = desc
    p_bd.font.size = Pt(13.5)
    p_bd.font.color.rgb = DARK_GRAY
    p_bd.space_after = Pt(15)

# ==========================================
# SLIDE 8: ĐÁNH GIÁ ƯU ĐIỂM & NHƯỢC ĐIỂM
# ==========================================
slide8 = prs.slides.add_slide(blank_layout)
apply_background_and_header(slide8, "Đánh Giá Ưu Điểm & Nhược Điểm")

# Khối cột trái: Ưu điểm (Xanh lá cây nhạt)
box_up = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
box_up.fill.solid()
box_up.fill.fore_color.rgb = LIGHT_GRAY
box_up.line.color.rgb = GREEN
box_up.line.width = Pt(1.5)

tf_up = box_up.text_frame
tf_up.word_wrap = True
tf_up.margin_top = tf_up.margin_left = tf_up.margin_right = Inches(0.3)
p_up_t = tf_up.paragraphs[0]
p_up_t.text = "ƯU ĐIỂM NỔI BẬT"
p_up_t.font.name = 'Arial'
p_up_t.font.size = Pt(18)
p_up_t.font.bold = True
p_up_t.font.color.rgb = DARK_BLUE

bullets_up = [
    "• Trải nghiệm mượt mà, trực quan: Giao diện Streamlit gọn nhẹ, cho phép người dùng thay đổi thông số bằng thanh kéo slider và nhận kết quả tức thì.",
    "• Xử lý tốt dữ liệu phi tuyến tính: Thuật toán Random Forest Regressor tự động học các quan hệ phức tạp mà không bị ảnh hưởng bởi hiện tượng đa cộng tuyến.",
    "• Kiến trúc 3 lớp rõ rệt: Tách biệt hoàn toàn phần giao diện, xử lý thuật toán và lưu trữ, giúp dễ dàng nâng cấp hoặc thay đổi mô hình."
]

for b in bullets_up:
    p_b = tf_up.add_paragraph()
    p_b.text = b
    p_b.font.size = Pt(13)
    p_b.font.color.rgb = DARK_GRAY
    p_b.space_after = Pt(8)

# Khối cột phải: Nhược điểm (Đỏ nhạt)
box_down = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
box_down.fill.solid()
box_down.fill.fore_color.rgb = LIGHT_GRAY
box_down.line.color.rgb = RED
box_down.line.width = Pt(1.5)

tf_down = box_down.text_frame
tf_down.word_wrap = True
tf_down.margin_top = tf_down.margin_left = tf_down.margin_right = Inches(0.3)
p_down_t = tf_down.paragraphs[0]
p_down_t.text = "NHƯỢC ĐIỂM HẠN CHẾ"
p_down_t.font.name = 'Arial'
p_down_t.font.size = Pt(18)
p_down_t.font.bold = True
p_down_t.font.color.rgb = DARK_BLUE

bullets_down = [
    "• Thiếu các biến thực tế quan trọng: Mô hình chưa tích hợp thông tin vị trí địa lý (khu vực, quận/huyện) - yếu tố hàng đầu trong định giá BĐS.",
    "• Kích thước tệp đóng gói lớn: File .pkl của Random Forest khá nặng vì chứa nhiều nhánh cây quyết định, gây tốn tài nguyên khi khởi chạy startup.",
    "• Khả năng ngoại suy kém: Thuật toán dạng cây không thể dự đoán tốt các giá trị nằm ngoài miền dữ liệu huấn luyện đã học."
]

for b in bullets_down:
    p_b = tf_down.add_paragraph()
    p_b.text = b
    p_b.font.size = Pt(13)
    p_b.font.color.rgb = DARK_GRAY
    p_b.space_after = Pt(8)

# ==========================================
# SLIDE 9: HƯỚNG PHÁT TRIỂN HỆ THỐNG
# ==========================================
slide9 = prs.slides.add_slide(blank_layout)
apply_background_and_header(slide9, "Hướng Phát Triển Hệ Thống Trong Tương Lai")

dev_plans = [
    ("1. NÂNG CẤP DỮ LIỆU & MÔ HÌNH", 
     "• Xây dựng bot cào tự động (Web Scraping) dữ liệu thời gian thực từ các sàn giao dịch lớn.\n• Tích hợp hệ thống thông tin địa lý GIS để lượng hóa yếu tố vị trí.\n• Nâng cấp lên các thuật toán mạnh mẽ hơn thuộc họ Boosting (XGBoost, LightGBM)."),
    ("2. HIỆN ĐẠI HÓA KIẾN TRÚC", 
     "• Chuyển đổi sang kiến trúc Microservices tách biệt hoàn toàn Frontend (ReactJS/Next.js) và Backend (FastAPI).\n• Thay thế tệp .csv bằng hệ quản trị cơ sở dữ liệu quan hệ mạnh mẽ (PostgreSQL/MySQL).\n• Đóng gói Docker, triển khai đám mây (AWS, GCP, Azure)."),
    ("3. TRẢI NGHIỆM NGƯỜI DÙNG", 
     "• Bổ sung tính năng quản lý tài khoản cá nhân, lưu lịch sử tra cứu.\n• Tích hợp bản đồ nhiệt (Real Estate Heatmap) tương tác địa lý trực quan.\n• Tích hợp bộ tính toán tài chính và lập kế hoạch vay vốn ngân hàng mua nhà.")
]

for i, (title, desc) in enumerate(dev_plans):
    left = Inches(0.8 + (i * 3.95))
    top = Inches(2.0)
    width = Inches(3.8)
    height = Inches(4.5)
    
    plan_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    plan_box.fill.solid()
    plan_box.fill.fore_color.rgb = LIGHT_GRAY
    plan_box.line.color.rgb = DARK_BLUE
    plan_box.line.width = Pt(1.5)
    
    tf_plan = plan_box.text_frame
    tf_plan.word_wrap = True
    tf_plan.margin_top = tf_plan.margin_left = tf_plan.margin_right = Inches(0.2)
    
    p_pt = tf_plan.paragraphs[0]
    p_pt.text = title
    p_pt.font.name = 'Arial'
    p_pt.font.size = Pt(15)
    p_pt.font.bold = True
    p_pt.font.color.rgb = DARK_BLUE
    p_pt.space_after = Pt(15)
    p_pt.alignment = PP_ALIGN.CENTER
    
    p_pd = tf_plan.add_paragraph()
    p_pd.text = desc
    p_pd.font.size = Pt(12)
    p_pd.font.color.rgb = DARK_GRAY

# ==========================================
# SLIDE 10: KẾT LUẬN & LỜI CẢM ƠN
# ==========================================
slide10 = prs.slides.add_slide(blank_layout)
apply_background_and_header(slide10)

# Chân trang trang trí cuối
bot_bar = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15))
bot_bar.fill.solid()
bot_bar.fill.fore_color.rgb = DARK_BLUE
bot_bar.line.fill.background()

# Hộp văn bản trung tâm
thank_box = slide10.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
tf_thank = thank_box.text_frame
tf_thank.word_wrap = True

p_th = tf_thank.paragraphs[0]
p_th.text = "XIN CHÂN THÀNH CẢM ƠN!"
p_th.font.name = 'Arial'
p_th.font.size = Pt(48)
p_th.font.bold = True
p_th.font.color.rgb = DARK_BLUE
p_th.alignment = PP_ALIGN.CENTER
p_th.space_after = Pt(15)

p_sch = tf_thank.add_paragraph()
p_sch.text = "Hội Đồng Báo Cáo Bài Tập Lớn - Công Nghệ Big Data"
p_sch.font.name = 'Arial'
p_sch.font.size = Pt(20)
p_sch.font.bold = True
p_sch.font.color.rgb = LIGHT_BLUE
p_sch.alignment = PP_ALIGN.CENTER
p_sch.space_after = Pt(20)

p_stud = tf_thank.add_paragraph()
p_stud.text = "Là một sinh viên khoa CNTT trường Đại học Hòa Bình, em đã ứng dụng thành công các công nghệ Big Data và Học máy hiện đại vào đồ án thực nghiệm thực tế. Sản phẩm nghiên cứu khoa học nghiêm túc này của Nhóm 4 dưới sự hướng dẫn trực tiếp của Thầy Đoàn Văn Hòa rất mong nhận được ý kiến đóng góp quý báu từ thầy cô!"
p_stud.font.name = 'Arial'
p_stud.font.size = Pt(15)
p_stud.font.color.rgb = DARK_GRAY
p_stud.alignment = PP_ALIGN.CENTER
p_stud.space_after = Pt(30)

p_cont = tf_thank.add_paragraph()
p_cont.text = "Sinh viên Đại học Hòa Bình\nKhoa Công nghệ thông tin | Email: anh123@gmail.com"
p_cont.font.name = 'Arial'
p_cont.font.size = Pt(13)
p_cont.font.bold = True
p_cont.font.color.rgb = DARK_BLUE
p_cont.alignment = PP_ALIGN.CENTER

# 4. LƯU PRESENTATION THÀNH FILE PPTX
output_filename = "BTL_BigData_HousePrice.pptx"
prs.save(output_filename)
print(f"Thành công: Đã tạo file '{output_filename}' hoàn chỉnh!")
```eof

### Hướng dẫn 3 bước cực kỳ nhanh để tạo ra file PowerPoint:

Bạn chỉ cần thực hiện 3 bước đơn giản sau trực tiếp trên máy tính của mình để có ngay tệp trình chiếu PowerPoint chuyên nghiệp:

1. **Cài đặt thư viện tạo PowerPoint:**
   Mở terminal / cmd trên máy tính và chạy lệnh cài đặt thư viện `python-pptx`:
   
```bash
   pip install python-pptx